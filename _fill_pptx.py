# -*- coding: utf-8 -*-
"""填充 Pi Agent 内容到 PPT 模板"""
import sys, io, copy
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

pptx_path = r"C:\newtask-pi\Agent Infra初赛方案PPT框架模板.pptx"
prs = Presentation(pptx_path)

# ── 颜色常量 ──
C_ORANGE = RGBColor(0xFF, 0x6B, 0x35)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK   = RGBColor(0x1B, 0x1F, 0x3B)
C_GRAY   = RGBColor(0x6B, 0x72, 0x80)
C_LIGHT  = RGBColor(0x9A, 0xA1, 0xB8)
C_BLUE   = RGBColor(0x3B, 0x82, 0xF6)
C_GREEN  = RGBColor(0x10, 0xB9, 0x81)
C_BG     = RGBColor(0x23, 0x28, 0x4A)

def set_text(shape, text, font_size=None, bold=None, color=None, alignment=None):
    """设置形状文本"""
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if alignment:
        p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = '微软雅黑'
    if font_size:
        run.font.size = Pt(font_size)
    if bold is not None:
        run.font.bold = bold
    if color:
        run.font.color.rgb = color

def set_multi_text(shape, lines, font_size=12, color=C_GRAY):
    """设置多行文本"""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = '微软雅黑'
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        p.space_after = Pt(3)

def replace_text_in_shape(shape, old_text, new_text):
    """在形状的文本中替换"""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            if old_text in run.text:
                run.text = run.text.replace(old_text, new_text)

def find_shape_by_text(slide, text):
    """通过文本查找形状"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            if text in shape.text_frame.text:
                return shape
    return None

def find_shapes_by_text(slide, texts):
    """通过多个文本查找形状列表"""
    results = []
    for text in texts:
        shape = find_shape_by_text(slide, text)
        if shape:
            results.append(shape)
    return results

slides = prs.slides

# ══════════════════════════════════════════════════════
# Slide 1: 封面
# ══════════════════════════════════════════════════════
s1 = slides[0]
for shape in s1.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if 'Agent Infra 新智基座' in txt:
            # 替换副标题
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = run.text.replace('Agent Infra 新智基座 初赛 · 方案 PPT 模板', 'Pi Agent — 多 Agent 协同智能助手平台')
        if '内容框架模板' in txt:
            tf = shape.text_frame
            tf.paragraphs[0].runs[0].text = '初赛方案 PPT'
            if len(tf.paragraphs) > 1:
                tf.paragraphs[1].runs[0].text = 'Pi Agent'
        if 'Datawhale · 参赛选手参考资料' in txt:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = 'Datawhale · Agent Infra 参赛作品'

print('✅ Slide 1: 封面 已填充')

# ══════════════════════════════════════════════════════
# Slide 2: P0 一页纸速览 — 6 个卡片
# ══════════════════════════════════════════════════════
s2 = slides[1]

# 卡片内容映射: (标题文本, 提示文本, 占位文本, 填充内容)
card_data = [
    ('项目名称', '≤ 20 字。一句话说清', '【在此填写项目名称】',
     'Pi Agent — 多 Agent 协同智能助手平台'),
    ('问题与场景', '目标用户是谁、在什么场景', '【描述真实场景与核心痛点】',
     '个人用户跨多信息源（网页/文件/股票/公众号）获取信息，现有聊天机器人无法多步推理、不能跨工具协同、缺乏长期记忆'),
    ('核心解决方案', '一句话讲清整体思路', '【概述端到端解决方案】',
     '双层 Agent Loop + 父Agent委托3类子Agent + 7个Skill + 19个工具 + 三层短期记忆 + 向量长期记忆'),
    ('创新点与差异化', '对比现有做法', '【列 1–2 个关键差异化优势】',
     '① 受控Agent运行时（/tasklist+质量门+修正）② 流式插话Steer ③ Skill文件化懒加载零代码扩展'),
    ('开放 / 复用价值', '可复用成果', '【说明复用与迁移价值】',
     'Skill以.md文件定义可零代码迁移；Agent Loop钩子化可复用任何LLM；记忆系统可独立部署'),
    ('当前进展', '做到了什么程度', '【说明当前完成度与里程碑】',
     '已完成可运行Web应用（FastAPI+React），支持多会话/流式输出/工具调用/子代理委托/长期记忆/股票分析'),
]

for title, hint, placeholder, content in card_data:
    # 找到占位文本形状并替换
    for shape in s2.shapes:
        if shape.has_text_frame and placeholder in shape.text_frame.text:
            set_multi_text(shape, [content], font_size=10, color=C_DARK)
            break

print('✅ Slide 2: 一页纸速览 已填充')

# ══════════════════════════════════════════════════════
# Slide 3: 目录 — 保持不变（已与模板一致）
# ══════════════════════════════════════════════════════
print('✅ Slide 3: 目录 保持不变')

# ══════════════════════════════════════════════════════
# Slide 4: 第一章标题 — 场景与价值
# ══════════════════════════════════════════════════════
print('✅ Slide 4: 第一章标题 保持不变')

# ══════════════════════════════════════════════════════
# Slide 5: 第一章内容 — 场景与价值
# ══════════════════════════════════════════════════════
s5 = slides[4]
# 替换 "建议覆盖..." 的提示文本
for shape in s5.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if '建议覆盖目标用户' in txt:
            lines = [
                '目标用户：知识工作者（研究员/分析师）、个人投资者、学生/学习者、内容创作者',
                '',
                '核心痛点：跨多个信息源手动切换工具效率低；查行情+看技术面+读公众号分析工具分散；缺乏统一入口和长期记忆',
                '',
                '真实场景：',
                '  ① 投资研究：用户问"看共进股份603118的技术面和ESG" → 自动获取K线→计算7因子→ESG评分→综合结论',
                '  ② 学习研究：用户问"读一下这个Datawhale教程" → web_fetch抓取→结构化摘要→学习路径建议',
                '  ③ 跨源整合：用户问"总结这篇公众号文章+搜相关资料" → wechat_article提取→委托research子代理搜索→综合输出',
                '',
                '可量化价值：跨源信息收集从15-30分钟降至30秒；股票技术分析全自动含ESG；三层记忆+向量长期记忆自动召回',
                '',
                '行业可复制性：金融（替换数据源用于基金/债券）、教育（Skill组合适配不同学科）、企业服务（Agent Loop用于运维排障）',
                '',
                '创新点：①受控Agent运行时 ②流式插话Steer ③Skill文件化懒加载',
            ]
            set_multi_text(shape, lines, font_size=11, color=C_DARK)

print('✅ Slide 5: 场景与价值 已填充')

# ══════════════════════════════════════════════════════
# Slide 6: 第二章标题 — 方案总览
# ══════════════════════════════════════════════════════
print('✅ Slide 6: 第二章标题 保持不变')

# ══════════════════════════════════════════════════════
# Slide 7: 第二章内容 — 方案总览（含架构图）
# ══════════════════════════════════════════════════════
s7 = slides[6]

# 替换 "建议用一张架构图..." 的提示文本为架构图内容
for shape in s7.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if '建议用一张架构图' in txt:
            lines = [
                '架构总览：',
                '  React Frontend（对话界面 | 股票分析面板 | 会话侧栏 | 富文本编辑器）',
                '    ↕ NDJSON Stream',
                '  FastAPI Backend',
                '    Chat Orchestrator（编排器）',
                '      ├─ Pi Agent 分支（/tasklist入口 + 质量门 + 自动修正）',
                '      └─ Agent Loop 双层循环（LLM ↔ 工具并发执行）',
                '           ├─ 主Agent（19个工具）',
                '           │    └─ 委托子代理：research(信息检索) / analysis(数据分析) / writer(内容创作)',
                '           ├─ Skill Registry（7个Skill: utility/reader/research/web/stock/wechat/webchat-art）',
                '           └─ Tool Registry（19个工具: calculator/web_search/stock_quote/github_repo/...）',
                '    ThreadState（短期记忆: recent+summary+pinned）',
                '    UserMemory（长期记忆: 向量搜索）',
                '    Steer Queue（流式插话）',
                '    ↕ DeepSeek API（OpenAI兼容）',
                '',
                '关键技术选型：DeepSeek API | FastAPI+asyncio | NDJSON流式协议 | DuckDB持久化 | React+TypeScript',
            ]
            set_multi_text(shape, lines, font_size=9, color=C_DARK)

print('✅ Slide 7: 方案总览 已填充')

# ══════════════════════════════════════════════════════
# Slide 8: 第三章标题 — 多Agent协同设计
# ══════════════════════════════════════════════════════
print('✅ Slide 8: 第三章标题 保持不变')

# ══════════════════════════════════════════════════════
# Slide 9: 第三章内容 — 多Agent协同设计
# ══════════════════════════════════════════════════════
s9 = slides[8]
for shape in s9.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if '建议覆盖 Agent 分工' in txt:
            lines = [
                'Agent 分工：',
                '  主Agent（父代理）：对话编排、任务分解、结果综合，拥有全部19个工具',
                '  Research子代理：信息检索专家，工具含 web_browse/local-text-read/get_weather 等',
                '  Analysis子代理：数据分析专家，工具含 calculator/unit_convert/text_transform',
                '  Writer子代理：内容创作专家，纯LLM无工具',
                '  Pi Agent：受控任务清单生成，显式入口+确定性质量门+自动修正',
                '',
                '任务拆解示例：用户"研究京东方，查行情+看技术面+搜新闻"',
                '  ├─ Task1: stock_quote 查实时行情（直接执行）',
                '  ├─ Task2: stock_analysis 技术分析（直接执行）',
                '  └─ Task3: 委托research子代理 → web_search+web_fetch → 返回结构化摘要',
                '  父Agent综合三个Task结果 → 输出最终回答',
                '',
                '上下文传递：父Agent→delegate_sub_agent(task,type)→创建子AgentContext(深度+1)→运行子agent_loop→事件转发到父流→tool_result返回',
                '',
                '异常处理：LLM超时3次重试+指数退避 | 工具超时30s/120s | 截断容错最多3次 | 子代理嵌套MAX_DEPTH=2',
                '',
                'Steer安全边界：只在turn边界生效（不中断工具执行）| 流结束返回409 | 以user角色注入',
            ]
            set_multi_text(shape, lines, font_size=10, color=C_DARK)

print('✅ Slide 9: 多Agent协同设计 已填充')

# ══════════════════════════════════════════════════════
# Slide 10: 第四章标题 — Skill工程体系
# ══════════════════════════════════════════════════════
print('✅ Slide 10: 第四章标题 保持不变')

# ══════════════════════════════════════════════════════
# Slide 11: 第四章内容 — Skill工程体系
# ══════════════════════════════════════════════════════
s11 = slides[10]
for shape in s11.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if '建议覆盖 Skill 清单' in txt:
            lines = [
                'Skill 清单（7个）：',
                '  utility-skill（实用工具）：数学计算/日期/单位换算 → 5个工具',
                '  reader-skill（信息读取）：文件读取/网页浏览/天气 → 5个工具',
                '  research-skill（研究助手）：子代理委托/复杂任务分解 → 5个工具',
                '  web-skill（网络研究）：搜索/抓取/GitHub/YouTube/PDF → 6个工具',
                '  stock-skill（股票行情）：A股实时行情/搜索/技术分析 → 4个工具',
                '  wechat-skill（微信公众号）：公众号文章提取分析 → 3个工具',
                '',
                'Skill 规格（front-matter + system_prompt）：',
                '  id | tool_names | output_policy | result_policy | routing_hints | fallback_policy',
                '  正文 = system_prompt（告诉LLM如何使用工具）',
                '',
                '失败处理：fallback_policy → direct-answer(无工具回答) / skip-capability(跳过) / retry(重试)',
                '',
                '生命周期：开发(.md文件) → 部署(discover扫描建索引) → 运行(首次get懒加载) → 更新(改文件重启)',
                '',
                '对官方Skills复用：兼容OpenAI Function Calling格式 | delegate_sub_agent可委托外部Agent | 可扩展MCP协议',
                '',
                '复用价值：Skill层(.md零代码迁移) | 工具层(独立函数拷贝+注册) | Agent Loop(钩子化可复用) | 记忆系统(独立部署)',
            ]
            set_multi_text(shape, lines, font_size=9, color=C_DARK)

print('✅ Slide 11: Skill工程体系 已填充')

# ══════════════════════════════════════════════════════
# Slide 12: 第五章标题 — 工程落地与安全可审计
# ══════════════════════════════════════════════════════
print('✅ Slide 12: 第五章标题 保持不变')

# ══════════════════════════════════════════════════════
# Slide 13: 第五章内容 — 工程落地与安全可审计
# ══════════════════════════════════════════════════════
s13 = slides[12]
for shape in s13.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if '建议覆盖可运行性' in txt:
            lines = [
                '可运行性：Web应用已部署（FastAPI+React），python app.py 即可启动，DuckDB零运维',
                '',
                '运行证据：',
                '  ① 股票分析：输入"000725" → 实时行情+7因子(MACD/RSI/KDJ/CCI/ROC/Williams/DMI)+ESG评分',
                '  ② 子代理委托："研究这篇公众号文章" → research子代理独立运行，前端显示层级和事件',
                '  ③ 流式插话：Agent执行中用户发Steer"换方向" → 下一turn调整路径',
                '',
                '可观测链路：',
                '  会话级: GET /api/conversations/{id}（ThreadState hydration）',
                '  长期记忆: GET /api/memories?session_id=xxx（向量搜索结果）',
                '  前端展示: 工具调用卡片 | 子代理层级 | Token用量统计',
                '  日志: NDJSON含 tool_call/tool_result/sub_agent_start/end 事件',
                '',
                '安全治理：XSS过滤(validators.py) | 参数化查询(防SQL注入) | 工具参数校验 | 子代理MAX_DEPTH=2',
                '  | 工具超时(30s/120s) | API Key不提交Git | LLM幻觉3次重试+回退',
                '',
                '云产品选型：DeepSeek API(必需) | DuckDB(嵌入式零运维) | Uvicorn(ASGI)',
                '  边界：不依赖外部向量DB/Redis/对象存储，保持轻量化',
            ]
            set_multi_text(shape, lines, font_size=9, color=C_DARK)

print('✅ Slide 13: 工程落地与安全 已填充')

# ══════════════════════════════════════════════════════
# Slide 14: 第六章标题 — 开放/开源计划
# ══════════════════════════════════════════════════════
print('✅ Slide 14: 第六章标题 保持不变')

# ══════════════════════════════════════════════════════
# Slide 15: 第六章内容 — 开放/开源计划
# ══════════════════════════════════════════════════════
s15 = slides[14]
for shape in s15.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if '建议覆盖可复用成果' in txt:
            lines = [
                '可复用成果：',
                '  Agent Loop：pi.dev风格双层循环，钩子化设计（拷贝agent_loop.py即用）',
                '  Skill系统：文件化定义，front-matter+system_prompt（拷贝skill_registry.py+skills/*.md）',
                '  记忆系统：三层短期+向量长期（拷贝thread_state.py+user_memory.py）',
                '  工具库：19个开箱即用工具（拷贝tools/目录）',
                '  子代理框架：支持递归委托、事件转发（拷贝sub_agent.py）',
                '  流式协议：NDJSON chunk+生命周期管理（拷贝stream/目录）',
                '',
                '接口契约：',
                '  工具注册：tool_registry.register(ChatToolDefinition(name, description, parameters, execute, ...))',
                '  Skill定义：skills/*.md（YAML front-matter + Markdown system_prompt）',
                '',
                '开源协议：MIT License',
                '第三方依赖：FastAPI / Pydantic / httpx / duckdb / OpenAI SDK（均MIT/Apache 2.0）',
            ]
            set_multi_text(shape, lines, font_size=9, color=C_DARK)

print('✅ Slide 15: 开源计划 已填充')

# ══════════════════════════════════════════════════════
# Slide 16: 第七章标题 — 落地计划与进展
# ══════════════════════════════════════════════════════
print('✅ Slide 16: 第七章标题 保持不变')

# ══════════════════════════════════════════════════════
# Slide 17: 第七章内容 — 落地计划与进展
# ══════════════════════════════════════════════════════
s17 = slides[16]
for shape in s17.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if '建议覆盖当前进展' in txt:
            lines = [
                '当前进展（全部 100% 完成）：',
                '  Agent Loop 双层循环 | 子代理系统(3类型) | Skill系统(7个) | 工具库(19个)',
                '  短期记忆(三层) | 长期记忆(向量) | Pi Agent受控运行时 | React前端 | 流式输出 | Steer插话',
                '',
                '里程碑：',
                '  MVP → 多Agent → Skill系统 → 记忆系统 → 流式插话 → 股票分析 → 公众号阅读 → Web研究',
                '',
                '落地计划：',
                '  初赛提交前：PPT完善 + Demo录制',
                '  决赛前：性能优化 + 30+工具 + 更多数据源',
                '  决赛阶段：对标AgentTeams，复刻运维故障排查场景',
                '  长期：开源社区，发布GitHub，贡献文档和教程',
                '',
                '风险控制：',
                '  LLM不稳定 → 3次重试+指数退避+回退直接回答',
                '  工具超时 → asyncio.wait_for(30s/120s)',
                '  子代理嵌套 → MAX_DEPTH=2',
                '  向量搜索失败 → 降级空数组，不影响主流程',
            ]
            set_multi_text(shape, lines, font_size=9, color=C_DARK)

print('✅ Slide 17: 落地计划与进展 已填充')

# ══════════════════════════════════════════════════════
# Slide 18: 第八章标题 — 团队介绍
# ══════════════════════════════════════════════════════
print('✅ Slide 18: 第八章标题 保持不变')

# ══════════════════════════════════════════════════════
# Slide 19: 第八章内容 — 团队介绍
# ══════════════════════════════════════════════════════
s19 = slides[18]
for shape in s19.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if '可从以下方面介绍团队' in txt or '基本背景信息' in txt:
            lines = [
                '1. 成员背景：',
                '   AI工程师 / 全栈开发者',
                '   核心技能：Python, FastAPI, React, LLM Agent系统, 向量数据库, TypeScript',
                '',
                '2. 团队分工：',
                '   架构设计：Agent Loop、Skill系统、记忆系统整体架构',
                '   后端开发：FastAPI、工具库(19个)、子代理、Pi Agent受控运行时',
                '   前端开发：React+TypeScript、流式解析、富文本编辑器、股票分析面板',
                '   测试优化：性能测试、错误处理、安全治理',
                '',
                '3. 团队成果：',
                '   已完成可运行的LLM Agent平台（FastAPI+React）',
                '   实现pi.dev风格Agent Loop双层循环',
                '   支持多Agent协同、流式插话、三层记忆+向量长期记忆',
                '   19个工具 + 7个Skill + 3类子代理',
            ]
            set_multi_text(shape, lines, font_size=10, color=C_DARK)

print('✅ Slide 19: 团队介绍 已填充')

# ── 保存 ──
output_path = r"C:\newtask-pi\Pi_Agent_初赛方案.pptx"
prs.save(output_path)
print(f'\n✅ PPT 已保存到: {output_path}')
