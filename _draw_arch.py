# -*- coding: utf-8 -*-
"""在 Slide 7 上绘制架构图（用形状模拟 SmartArt 效果）"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

pptx_path = r"C:\newtask-pi\Pi_Agent_初赛方案.pptx"
prs = Presentation(pptx_path)

# ── 颜色 ──
C_ORANGE = RGBColor(0xFF, 0x6B, 0x35)
C_ORANGE_L = RGBColor(0xFF, 0xE5, 0xD9)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK   = RGBColor(0x1B, 0x1F, 0x3B)
C_DARK_L = RGBColor(0x2D, 0x33, 0x5C)
C_GRAY   = RGBColor(0x6B, 0x72, 0x80)
C_BLUE   = RGBColor(0x3B, 0x82, 0xF6)
C_BLUE_L = RGBColor(0xDB, 0xEA, 0xFE)
C_GREEN  = RGBColor(0x10, 0xB9, 0x81)
C_GREEN_L= RGBColor(0xD1, 0xFA, 0xE5)
C_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
C_PURPLE_L=RGBColor(0xED, 0xE9, 0xFE)
C_TEAL   = RGBColor(0x14, 0xB8, 0xA6)
C_TEAL_L = RGBColor(0xCC, 0xFB, 0xF1)
C_RED    = RGBColor(0xEF, 0x44, 0x44)
C_RED_L  = RGBColor(0xFE, 0xE2, 0xE2)
C_BG     = RGBColor(0xF8, 0xFA, 0xFC)
C_LINE   = RGBColor(0xCB, 0xD5, 0xE1)

def add_box(slide, left, top, width, height, text, fill_color, text_color=C_WHITE, font_size=10, bold=True, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """添加一个带文本的彩色框"""
    shape = slide.shapes.add_shape(shape_type, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # 无边框
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.name = '微软雅黑'
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    # 垂直居中
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # 设置内边距
    tf.margin_top = Emu(36000)
    tf.margin_bottom = Emu(36000)
    tf.margin_left = Emu(72000)
    tf.margin_right = Emu(72000)
    return shape

def add_arrow(slide, left, top, width, height, color=C_GRAY, direction='down'):
    """添加箭头"""
    if direction == 'down':
        shape_type = MSO_SHAPE.DOWN_ARROW
    elif direction == 'right':
        shape_type = MSO_SHAPE.RIGHT_ARROW
    elif direction == 'up':
        shape_type = MSO_SHAPE.UP_ARROW
    else:
        shape_type = MSO_SHAPE.DOWN_ARROW
    shape = slide.shapes.add_shape(shape_type, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_label(slide, left, top, width, height, text, color=C_GRAY, font_size=8, bold=False):
    """添加纯文本标签"""
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = '微软雅黑'
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

# ══════════════════════════════════════════════════════
# Slide 7: 绘制架构图
# ══════════════════════════════════════════════════════
s7 = prs.slides[6]

# 先删除之前填入的文本内容形状（找到含"架构总览"的文本框）
shapes_to_remove = []
for shape in s7.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if '架构总览' in txt or 'React Frontend' in txt:
            shapes_to_remove.append(shape)

for shape in shapes_to_remove:
    sp = shape._element
    sp.getparent().remove(sp)

print(f'已删除 {len(shapes_to_remove)} 个旧文本框')

# ── 布局参数 ──
# Slide: 12192000 x 6858000 EMU
# 内容区: left=594360, top=1297920, w=10972800, h=5335680 (大致)
LEFT = 594360
TOP = 1250000
CW = 10972800  # 内容宽度

# ─── Layer 1: 前端 ───
y1 = TOP
h1 = 420000
# 标签
add_text_label(s7, LEFT, y1 - 200000, 800000, 200000, '前端层', C_ORANGE, font_size=9, bold=True)
# 4个前端模块
mod_w = 2500000
gap = (CW - mod_w * 4) // 3
fe_modules = [
    ('React 对话界面', C_BLUE),
    ('股票分析面板', C_BLUE),
    ('会话侧栏', C_BLUE),
    ('富文本编辑器', C_BLUE),
]
for i, (text, color) in enumerate(fe_modules):
    x = LEFT + i * (mod_w + gap)
    add_box(s7, x, y1, mod_w, h1, text, color, font_size=10)

# 箭头: 前端 → 后端
arrow_y = y1 + h1 + 30000
add_arrow(s7, LEFT + CW//2 - 150000, arrow_y, 300000, 50000, C_GRAY, 'down')
add_text_label(s7, LEFT + CW//2 + 200000, arrow_y - 5000, 1500000, 50000, 'NDJSON Stream', C_GRAY, font_size=8)

# ─── Layer 2: 后端核心 ───
y2 = arrow_y + 80000
h2 = 420000
add_text_label(s7, LEFT, y2 - 200000, 800000, 200000, '后端层', C_ORANGE, font_size=9, bold=True)
# Chat Orchestrator
add_box(s7, LEFT, y2, CW, h2, 'FastAPI Backend — Chat Orchestrator（编排器：会话管理 / 流式分发 / 工具调用 / 子代理委托）', C_DARK, font_size=11)

# 箭头
arrow_y2 = y2 + h2 + 30000
add_arrow(s7, LEFT + CW//2 - 150000, arrow_y2, 300000, 50000, C_GRAY, 'down')

# ─── Layer 3: Agent Loop (核心) ───
y3 = arrow_y2 + 80000
h3 = 460000
add_text_label(s7, LEFT, y3 - 200000, 1000000, 200000, 'Agent Loop（双层循环）', C_ORANGE, font_size=9, bold=True)
# Pi Agent 分支
pi_w = 3200000
pi_x = LEFT
add_box(s7, pi_x, y3, pi_w, h3, 'Pi Agent 受控运行时\n/tasklist → 质量门 → 自动修正', C_ORANGE, font_size=9)
# Agent Loop 核心
al_w = 3200000
al_x = LEFT + pi_w + gap
add_box(s7, al_x, y3, al_w, h3, 'Agent Loop 双层循环\nLLM ↔ 工具并发执行（asyncio）', C_PURPLE, font_size=9)
# Steer
st_w = CW - pi_w - al_w - gap * 2
st_x = LEFT + pi_w + al_w + gap * 2
add_box(s7, st_x, y3, st_w, h3, 'Steer 流式插话\n用户中途干预', C_RED, font_size=9)

# 箭头
arrow_y3 = y3 + h3 + 30000
add_arrow(s7, LEFT + CW//2 - 150000, arrow_y3, 300000, 50000, C_GRAY, 'down')

# ─── Layer 4: 子代理 + Skill + 工具 ───
y4 = arrow_y3 + 80000
h4 = 420000
add_text_label(s7, LEFT, y4 - 200000, 1200000, 200000, '子代理 / Skill / 工具', C_ORANGE, font_size=9, bold=True)

# 3个子代理 + 1个Skill + 1个工具
sub_w = 2000000
sub_gap = (CW - sub_w * 5) // 4

# 子代理1: Research
sub_items = [
    ('Research\n子代理', C_GREEN, '信息检索\n5个工具'),
    ('Analysis\n子代理', C_TEAL, '数据分析\n3个工具'),
    ('Writer\n子代理', C_BLUE, '内容创作\n纯LLM'),
    ('Skill 系统\n(7个 .md)', C_ORANGE, 'utility/reader\nresearch/web\nstock/wechat'),
    ('Tool Registry\n(19个工具)', C_PURPLE, 'calculator\nweb_search\nstock_quote\nwechat_article'),
]
for i, (text, color, desc) in enumerate(sub_items):
    x = LEFT + i * (sub_w + sub_gap)
    # 主框
    add_box(s7, x, y4, sub_w, 280000, text, color, font_size=9)
    # 描述框
    add_box(s7, x, y4 + 290000, sub_w, 130000, desc, C_WHITE, text_color=C_GRAY, font_size=7, bold=False, shape_type=MSO_SHAPE.RECTANGLE)

# ─── Layer 5: 记忆 + LLM ───
y5 = y4 + h4 + 250000
h5 = 380000
add_text_label(s7, LEFT, y5 - 200000, 800000, 200000, '记忆 / LLM', C_ORANGE, font_size=9, bold=True)

# 3个记忆模块 + 1个LLM
mem_w = 2500000
mem_gap = (CW - mem_w * 4) // 3
mem_items = [
    ('ThreadState\n三层短期记忆\n(recent+summary+pinned)', C_TEAL),
    ('UserMemory\n长期记忆\n(向量搜索 DuckDB)', C_GREEN),
    ('DeepSeek API\n(LLM 推理引擎)', C_DARK),
    ('DuckDB\n会话持久化\n(嵌入式零运维)', C_GRAY),
]
for i, (text, color) in enumerate(mem_items):
    x = LEFT + i * (mem_w + mem_gap)
    add_box(s7, x, y5, mem_w, h5, text, color, font_size=8)

print('✅ 架构图已绘制')

# ── 保存 ──
output_path = r"C:\newtask-pi\Pi_Agent_初赛方案.pptx"
prs.save(output_path)
print(f'\n✅ PPT 已保存: {output_path}')
