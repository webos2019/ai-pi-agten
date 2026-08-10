import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation

pptx_path = r"C:\newtask-pi\Agent Infra初赛方案PPT框架模板.pptx"

prs = Presentation(pptx_path)
print(f"幻灯片总数: {len(prs.slides)}")
print("=" * 70)

for i, slide in enumerate(prs.slides, 1):
    print(f"\n{'='*70}")
    print(f"  第 {i} 页幻灯片")
    print(f"{'='*70}")

    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(f"  {text}")

        if shape.has_table:
            print("  [表格]")
            table = shape.table
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                print(f"    | {'  |  '.join(cells)} |")

    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            print(f"  [备注] {notes}")
