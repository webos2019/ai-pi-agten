import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

pptx_path = r"C:\newtask-pi\Agent Infra初赛方案PPT框架模板.pptx"
prs = Presentation(pptx_path)

print(f"Slide width: {prs.slide_width}, height: {prs.slide_height}")
print(f"Slides: {len(prs.slides)}")
print()

for i, slide in enumerate(prs.slides, 1):
    print(f"=== Slide {i} ===")
    print(f"  Layout: {slide.slide_layout.name}")
    for j, shape in enumerate(slide.shapes):
        print(f"  Shape {j}: name={shape.name}, type={shape.shape_type}")
        print(f"    pos: left={shape.left}, top={shape.top}, w={shape.width}, h={shape.height}")
        if shape.has_text_frame:
            for k, para in enumerate(shape.text_frame.paragraphs):
                text = para.text.strip()
                if text:
                    print(f"    Para {k}: '{text[:80]}'")
                    for run in para.runs:
                        print(f"      Run: font={run.font.name}, size={run.font.size}, bold={run.font.bold}, color={run.font.color.rgb if run.font.color and run.font.color.type else 'None'}")
    print()
